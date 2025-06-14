from mcp.server.fastmcp import FastMCP
from dotenv import load_dotenv
import os
import re
import json
import sys
from typing import List, Dict, Any, Optional, Tuple
from collections import Counter
import base64
from datetime import datetime
from pathlib import Path

# PowerPoint generation
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.dml import MSO_THEME_COLOR

# Image processing
from PIL import Image
import io

# PDF analysis integration
import PyPDF2
import fitz  # PyMuPDF
from sklearn.feature_extraction.text import TfidfVectorizer
import numpy as np

print("Starting PowerPoint Generator MCP Server...", file=sys.stderr)

# Load environment variables
try:
    load_dotenv("../.env")
    print("Environment loaded", file=sys.stderr)
except Exception as e:
    print(f"Environment load error: {e}", file=sys.stderr)

# Initialize MCP
mcp = FastMCP("ppt-generator")


class BrandGuidelines:
    """Brand identity and styling guidelines"""

    MCKINSEY_STYLE = {
        "primary_color": (0, 47, 108),  # McKinsey Blue
        "secondary_color": (255, 255, 255),  # White
        "accent_color": (220, 53, 69),  # Red for highlights
        "text_color": (51, 51, 51),  # Dark Gray
        "light_gray": (128, 128, 128),
        "font_title": "Calibri",
        "font_body": "Calibri",
        "font_size_title": 28,
        "font_size_subtitle": 20,
        "font_size_body": 14,
        "font_size_caption": 10,
        "slide_margin": 0.5,
        "bullet_style": "•",
        "chart_colors": [
            (0, 47, 108),
            (220, 53, 69),
            (40, 167, 69),
            (255, 193, 7),
            (108, 117, 125),
        ],
    }

    CORPORATE_STYLE = {
        "primary_color": (25, 25, 112),  # Midnight Blue
        "secondary_color": (245, 245, 245),  # Light Gray
        "accent_color": (255, 165, 0),  # Orange
        "text_color": (0, 0, 0),  # Black
        "light_gray": (169, 169, 169),
        "font_title": "Arial",
        "font_body": "Arial",
        "font_size_title": 24,
        "font_size_subtitle": 18,
        "font_size_body": 12,
        "font_size_caption": 9,
        "slide_margin": 0.4,
        "bullet_style": "►",
        "chart_colors": [
            (25, 25, 112),
            (255, 165, 0),
            (50, 205, 50),
            (255, 69, 0),
            (138, 43, 226),
        ],
    }

    MODERN_STYLE = {
        "primary_color": (72, 61, 139),  # Dark Slate Blue
        "secondary_color": (248, 249, 250),  # Off White
        "accent_color": (220, 20, 60),  # Crimson
        "text_color": (33, 37, 41),  # Dark
        "light_gray": (108, 117, 125),
        "font_title": "Segoe UI",
        "font_body": "Segoe UI",
        "font_size_title": 26,
        "font_size_subtitle": 19,
        "font_size_body": 13,
        "font_size_caption": 10,
        "slide_margin": 0.3,
        "bullet_style": "▶",
        "chart_colors": [
            (72, 61, 139),
            (220, 20, 60),
            (32, 201, 151),
            (253, 126, 20),
            (156, 39, 176),
        ],
    }


class SlideLayoutEngine:
    """Intelligent slide layout and hierarchy engine"""

    def __init__(self, brand_style: Dict):
        self.brand_style = brand_style

    def determine_slide_type(self, content: Dict) -> str:
        """Determine the best slide type based on content"""
        if content.get("type") == "title":
            return "title"
        elif content.get("charts") or content.get("data"):
            return "chart"
        elif len(content.get("bullets", [])) > 6:
            return "content_heavy"
        elif content.get("image"):
            return "image_content"
        elif content.get("comparison"):
            return "comparison"
        else:
            return "standard"

    def get_layout_template(self, slide_type: str) -> Dict:
        """Get layout specifications for different slide types"""
        layouts = {
            "title": {
                "title_size": self.brand_style["font_size_title"] + 4,
                "subtitle_size": self.brand_style["font_size_subtitle"],
                "title_position": (1, 2, 8, 1.5),
                "subtitle_position": (1, 4, 8, 1),
                "content_position": None,
            },
            "standard": {
                "title_size": self.brand_style["font_size_title"],
                "title_position": (0.5, 0.5, 9, 1),
                "content_position": (0.5, 2, 9, 5.5),
            },
            "chart": {
                "title_size": self.brand_style["font_size_title"],
                "title_position": (0.5, 0.3, 9, 0.8),
                "chart_position": (1, 1.5, 8, 5),
                "caption_position": (1, 6.8, 8, 0.5),
            },
            "content_heavy": {
                "title_size": self.brand_style["font_size_title"] - 2,
                "title_position": (0.5, 0.3, 9, 0.7),
                "content_position": (0.5, 1.2, 4.4, 6),
                "content_position_2": (5, 1.2, 4.4, 6),
            },
            "image_content": {
                "title_size": self.brand_style["font_size_title"],
                "title_position": (0.5, 0.3, 9, 0.8),
                "image_position": (0.5, 1.5, 4.5, 4),
                "content_position": (5.2, 1.5, 4.3, 4),
            },
            "comparison": {
                "title_size": self.brand_style["font_size_title"],
                "title_position": (0.5, 0.3, 9, 0.8),
                "left_content_position": (0.5, 1.5, 4.2, 5),
                "right_content_position": (5.3, 1.5, 4.2, 5),
            },
        }
        return layouts.get(slide_type, layouts["standard"])


class ContentSafeguards:
    """Safeguards for brand voice, confidentiality, and content quality"""

    CONFIDENTIAL_PATTERNS = [
        r"\b(?:confidential|proprietary|internal only|restricted)\b",
        r"\b(?:salary|compensation|budget|cost|revenue|profit)\s*[:=]\s*[\d,.$]+",
        r"\b(?:ssn|social security|credit card|account number)\b",
        r"\b[A-Z]{2,}\s*[#-]?\s*\d{4,}\b",  # Possible ID numbers
    ]

    INAPPROPRIATE_CONTENT = [
        r"\b(?:confidential|secret|proprietary|internal)\s+(?:information|data|document)\b",
        r"\b(?:do not distribute|not for distribution|internal use only)\b",
    ]

    @staticmethod
    def check_content_safety(text: str) -> Dict[str, Any]:
        """Check content for confidentiality and appropriateness issues"""
        issues = []

        # Check for confidential patterns
        for pattern in ContentSafeguards.CONFIDENTIAL_PATTERNS:
            matches = re.findall(pattern, text, re.IGNORECASE)
            if matches:
                issues.append(
                    {
                        "type": "confidential",
                        "pattern": pattern,
                        "matches": matches[:3],  # Limit to first 3 matches
                    }
                )

        # Check for inappropriate content
        for pattern in ContentSafeguards.INAPPROPRIATE_CONTENT:
            matches = re.findall(pattern, text, re.IGNORECASE)
            if matches:
                issues.append(
                    {
                        "type": "inappropriate",
                        "pattern": pattern,
                        "matches": matches[:3],
                    }
                )

        return {
            "safe": len(issues) == 0,
            "issues": issues,
            "risk_level": "high"
            if len(issues) > 2
            else "medium"
            if len(issues) > 0
            else "low",
        }

    @staticmethod
    def sanitize_content(text: str) -> str:
        """Remove or mask potentially sensitive content"""
        # Mask potential sensitive numbers
        text = re.sub(r"\b\d{4,}\b", "[REDACTED]", text)

        # Remove confidential markers
        text = re.sub(
            r"\b(?:confidential|proprietary|internal only)\s*[:.-]?\s*",
            "",
            text,
            flags=re.IGNORECASE,
        )

        return text


class PowerPointGenerator:
    """Main PowerPoint generation engine"""

    def __init__(self, brand_style: str = "mckinsey"):
        self.brand_styles = {
            "mckinsey": BrandGuidelines.MCKINSEY_STYLE,
            "corporate": BrandGuidelines.CORPORATE_STYLE,
            "modern": BrandGuidelines.MODERN_STYLE,
        }
        self.brand_style = self.brand_styles.get(
            brand_style, BrandGuidelines.MCKINSEY_STYLE
        )
        self.layout_engine = SlideLayoutEngine(self.brand_style)

    def create_presentation(self) -> Presentation:
        """Create a new presentation with brand styling"""
        prs = Presentation()

        # Set slide size (16:9 aspect ratio)
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        return prs

    def add_title_slide(
        self, prs: Presentation, title: str, subtitle: str = "", author: str = ""
    ) -> None:
        """Add a professionally styled title slide"""
        slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Configure title
        title_shape = slide.shapes.title
        title_shape.text = title
        title_paragraph = title_shape.text_frame.paragraphs[0]
        title_paragraph.font.name = self.brand_style["font_title"]
        title_paragraph.font.size = Pt(self.brand_style["font_size_title"] + 6)
        title_paragraph.font.color.rgb = RGBColor(*self.brand_style["primary_color"])
        title_paragraph.alignment = PP_ALIGN.CENTER

        # Configure subtitle
        if slide.shapes.placeholders.count > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_text = f"{subtitle}\n\n{author}" if author else subtitle
            subtitle_shape.text = subtitle_text
            subtitle_paragraph = subtitle_shape.text_frame.paragraphs[0]
            subtitle_paragraph.font.name = self.brand_style["font_body"]
            subtitle_paragraph.font.size = Pt(self.brand_style["font_size_subtitle"])
            subtitle_paragraph.font.color.rgb = RGBColor(
                *self.brand_style["text_color"]
            )
            subtitle_paragraph.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, prs: Presentation, slide_content: Dict) -> None:
        """Add a content slide with intelligent layout"""
        slide_type = self.layout_engine.determine_slide_type(slide_content)
        layout_spec = self.layout_engine.get_layout_template(slide_type)

        slide_layout = prs.slide_layouts[1]  # Content slide layout
        slide = prs.slides.add_slide(slide_layout)

        # Add title
        title_shape = slide.shapes.title
        title_shape.text = slide_content.get("title", "Untitled Slide")
        title_paragraph = title_shape.text_frame.paragraphs[0]
        self._style_text(title_paragraph, "title")

        # Add content based on slide type
        if slide_type == "chart" and slide_content.get("chart_data"):
            self._add_chart(slide, slide_content["chart_data"], layout_spec)
        elif slide_type == "content_heavy":
            self._add_two_column_content(slide, slide_content, layout_spec)
        elif slide_type == "comparison":
            self._add_comparison_content(slide, slide_content, layout_spec)
        else:
            self._add_standard_content(slide, slide_content, layout_spec)

    def _style_text(self, paragraph, text_type: str):
        """Apply brand styling to text"""
        font_sizes = {
            "title": self.brand_style["font_size_title"],
            "subtitle": self.brand_style["font_size_subtitle"],
            "body": self.brand_style["font_size_body"],
            "caption": self.brand_style["font_size_caption"],
        }

        paragraph.font.name = (
            self.brand_style["font_title"]
            if text_type == "title"
            else self.brand_style["font_body"]
        )
        paragraph.font.size = Pt(
            font_sizes.get(text_type, self.brand_style["font_size_body"])
        )

        if text_type == "title":
            paragraph.font.color.rgb = RGBColor(*self.brand_style["primary_color"])
            paragraph.font.bold = True
        else:
            paragraph.font.color.rgb = RGBColor(*self.brand_style["text_color"])

    def _add_standard_content(self, slide, content: Dict, layout_spec: Dict):
        """Add standard bullet-point content"""
        if not content.get("bullets"):
            return

        # Create text box
        left, top, width, height = [Inches(x) for x in layout_spec["content_position"]]
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame
        text_frame.word_wrap = True

        # Add bullets
        for i, bullet in enumerate(content["bullets"][:8]):  # Limit bullets
            if i == 0:
                p = text_frame.paragraphs[0]
            else:
                p = text_frame.add_paragraph()

            p.text = bullet
            p.level = 0
            self._style_text(p, "body")

    def _add_two_column_content(self, slide, content: Dict, layout_spec: Dict):
        """Add content in two columns for content-heavy slides"""
        bullets = content.get("bullets", [])
        mid_point = len(bullets) // 2

        # Left column
        left, top, width, height = [Inches(x) for x in layout_spec["content_position"]]
        left_textbox = slide.shapes.add_textbox(left, top, width, height)
        left_frame = left_textbox.text_frame

        for i, bullet in enumerate(bullets[:mid_point]):
            if i == 0:
                p = left_frame.paragraphs[0]
            else:
                p = left_frame.add_paragraph()
            p.text = bullet
            self._style_text(p, "body")

        # Right column
        left, top, width, height = [
            Inches(x) for x in layout_spec["content_position_2"]
        ]
        right_textbox = slide.shapes.add_textbox(left, top, width, height)
        right_frame = right_textbox.text_frame

        for i, bullet in enumerate(bullets[mid_point:]):
            if i == 0:
                p = right_frame.paragraphs[0]
            else:
                p = right_frame.add_paragraph()
            p.text = bullet
            self._style_text(p, "body")

    def _add_comparison_content(self, slide, content: Dict, layout_spec: Dict):
        """Add comparison content with left/right layout"""
        comparison = content.get("comparison", {})

        # Left side
        left, top, width, height = [
            Inches(x) for x in layout_spec["left_content_position"]
        ]
        left_textbox = slide.shapes.add_textbox(left, top, width, height)
        left_frame = left_textbox.text_frame

        # Add left title
        p = left_frame.paragraphs[0]
        p.text = comparison.get("left_title", "Option A")
        self._style_text(p, "subtitle")
        p.font.bold = True

        # Add left bullets
        for bullet in comparison.get("left_points", []):
            p = left_frame.add_paragraph()
            p.text = bullet
            self._style_text(p, "body")

        # Right side
        left, top, width, height = [
            Inches(x) for x in layout_spec["right_content_position"]
        ]
        right_textbox = slide.shapes.add_textbox(left, top, width, height)
        right_frame = right_textbox.text_frame

        # Add right title
        p = right_frame.paragraphs[0]
        p.text = comparison.get("right_title", "Option B")
        self._style_text(p, "subtitle")
        p.font.bold = True

        # Add right bullets
        for bullet in comparison.get("right_points", []):
            p = right_frame.add_paragraph()
            p.text = bullet
            self._style_text(p, "body")

    def _add_chart(self, slide, chart_data: Dict, layout_spec: Dict):
        """Add a chart to the slide"""
        try:
            chart_type = chart_data.get("type", "column")

            # Create chart data
            chart_data_obj = CategoryChartData()
            chart_data_obj.categories = chart_data.get(
                "categories", ["Category 1", "Category 2"]
            )
            chart_data_obj.add_series("Series 1", chart_data.get("values", [1, 2]))

            # Add chart
            left, top, width, height = [
                Inches(x) for x in layout_spec["chart_position"]
            ]

            chart_type_map = {
                "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
                "bar": XL_CHART_TYPE.BAR_CLUSTERED,
                "line": XL_CHART_TYPE.LINE,
                "pie": XL_CHART_TYPE.PIE,
            }

            chart = slide.shapes.add_chart(
                chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED),
                left,
                top,
                width,
                height,
                chart_data_obj,
            ).chart

            # Style chart
            chart.has_legend = True
            chart.legend.position = 2  # Right side

        except Exception as e:
            print(f"Chart creation failed: {e}", file=sys.stderr)
            # Fallback to text representation
            self._add_chart_fallback(slide, chart_data, layout_spec)

    def _add_chart_fallback(self, slide, chart_data: Dict, layout_spec: Dict):
        """Add chart data as formatted text when chart creation fails"""
        left, top, width, height = [Inches(x) for x in layout_spec["chart_position"]]
        textbox = slide.shapes.add_textbox(left, top, width, height)
        text_frame = textbox.text_frame

        p = text_frame.paragraphs[0]
        p.text = f"Chart: {chart_data.get('title', 'Data Visualization')}"
        self._style_text(p, "subtitle")

        categories = chart_data.get("categories", [])
        values = chart_data.get("values", [])

        for cat, val in zip(categories, values):
            p = text_frame.add_paragraph()
            p.text = f"• {cat}: {val}"
            self._style_text(p, "body")


class ContentProcessor:
    """Process and structure content for presentation generation"""

    @staticmethod
    def extract_structure_from_pdf(pdf_path: str) -> Dict[str, Any]:
        """Extract structured content from PDF for presentation"""
        try:
            # Use the PDF analyzer functionality
            doc = fitz.open(pdf_path)
            text = ""
            for page in doc:
                text += page.get_text()
            doc.close()

            # Basic structure extraction
            sections = ContentProcessor._identify_sections(text)
            key_points = ContentProcessor._extract_key_points(text)

            return {
                "sections": sections,
                "key_points": key_points,
                "source_file": pdf_path,
            }
        except Exception as e:
            return {"error": f"PDF processing failed: {str(e)}"}

    @staticmethod
    def _identify_sections(text: str) -> List[Dict]:
        """Identify sections and headers in text"""
        lines = text.split("\n")
        sections = []
        current_section = None

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Check if line looks like a header
            if len(line) < 100 and (
                line.isupper()
                or re.match(r"^\d+\.?\s+[A-Z]", line)
                or line.endswith(":")
            ):
                if current_section:
                    sections.append(current_section)

                current_section = {"title": line.rstrip(":"), "content": []}
            elif current_section:
                current_section["content"].append(line)

        if current_section:
            sections.append(current_section)

        return sections

    @staticmethod
    def _extract_key_points(text: str) -> List[str]:
        """Extract key bullet points from text"""
        # Look for bullet-like patterns
        bullet_patterns = [
            r"^\s*[•▪▫‣⁃]\s+(.+)",
            r"^\s*[-*]\s+(.+)",
            r"^\s*\d+[.)]\s+(.+)",
            r"^\s*[a-z][.)]\s+(.+)",
        ]

        key_points = []
        lines = text.split("\n")

        for line in lines:
            line = line.strip()
            for pattern in bullet_patterns:
                match = re.match(pattern, line, re.IGNORECASE)
                if match:
                    point = match.group(1).strip()
                    if len(point) > 10 and len(point) < 200:  # Reasonable length
                        key_points.append(point)
                    break

        return key_points[:20]  # Limit to top 20 points


# MCP Tools


@mcp.tool()
async def generate_presentation(
    title: str,
    content_source: str = "",
    pdf_path: str = "",
    brand_style: str = "mckinsey",
    slide_count: int = 10,
    include_charts: bool = True,
    output_path: str = "",
    presentation_type: str = "proposal",
) -> str:
    """
    Generate a professional PowerPoint presentation with McKinsey-style formatting.

    Args:
        title: Presentation title
        content_source: Text content or topic for the presentation
        pdf_path: Path to PDF file to extract content from
        brand_style: Style template (mckinsey, corporate, modern)
        slide_count: Target number of slides
        include_charts: Whether to include data visualizations
        output_path: Where to save the presentation
        presentation_type: Type of presentation (proposal, brand-identity, executive-summary)

    Returns:
        JSON string with generation results and file path
    """
    try:
        print(f"Generating presentation: {title}", file=sys.stderr)

        # Initialize generator
        generator = PowerPointGenerator(brand_style)
        prs = generator.create_presentation()

        # Process content source
        if pdf_path and os.path.exists(pdf_path):
            pdf_content = ContentProcessor.extract_structure_from_pdf(pdf_path)
            if "error" in pdf_content:
                return json.dumps({"error": pdf_content["error"]})
        else:
            pdf_content = None

        # Check content safety
        safety_check = ContentSafeguards.check_content_safety(content_source)
        if not safety_check["safe"]:
            print(
                f"Content safety issues detected: {safety_check['issues']}",
                file=sys.stderr,
            )
            content_source = ContentSafeguards.sanitize_content(content_source)

        # Generate presentation structure based on type
        slide_structure = _generate_slide_structure(
            presentation_type, title, content_source, pdf_content, slide_count
        )

        # Add title slide
        generator.add_title_slide(
            prs,
            title,
            slide_structure.get("subtitle", ""),
            slide_structure.get("author", "Executive Presentation"),
        )

        # Add content slides
        for slide_content in slide_structure["slides"]:
            generator.add_content_slide(prs, slide_content)

        # Save presentation
        if not output_path:
            timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
            filename = f"{title.replace(' ', '_')}_{timestamp}.pptx"
            output_path = os.path.join(os.getcwd(), filename)

        prs.save(output_path)

        result = {
            "success": True,
            "file_path": output_path,
            "slide_count": len(prs.slides),
            "brand_style": brand_style,
            "presentation_type": presentation_type,
            "safety_check": safety_check,
            "generated_at": datetime.now().isoformat(),
        }

        print(f"Presentation saved: {output_path}", file=sys.stderr)
        return json.dumps(result, indent=2)

    except Exception as e:
        error_msg = f"Presentation generation failed: {str(e)}"
        print(error_msg, file=sys.stderr)
        return json.dumps({"error": error_msg})


@mcp.tool()
async def analyze_slide_template(image_path: str) -> str:
    """
    Analyze an uploaded slide template image to extract design elements and structure.

    Args:
        image_path: Path to the slide template image

    Returns:
        JSON string with analyzed design elements and recommendations
    """
    try:
        if not os.path.exists(image_path):
            return json.dumps({"error": f"Image file not found: {image_path}"})

        # Basic image analysis
        with Image.open(image_path) as img:
            width, height = img.size
            aspect_ratio = width / height

            # Convert to RGB if needed
            if img.mode != "RGB":
                img = img.convert("RGB")

            # Simple color analysis - get dominant colors
            colors = img.getcolors(maxcolors=256 * 256 * 256)
            if colors:
                dominant_colors = sorted(colors, key=lambda x: x[0], reverse=True)[:5]
                color_palette = [color[1] for color in dominant_colors]
            else:
                color_palette = [(255, 255, 255)]  # Default to white

        analysis_result = {
            "image_path": image_path,
            "dimensions": {"width": width, "height": height},
            "aspect_ratio": round(aspect_ratio, 2),
            "slide_format": "16:9"
            if 1.7 <= aspect_ratio <= 1.8
            else "4:3"
            if 1.3 <= aspect_ratio <= 1.4
            else "custom",
            "dominant_colors": color_palette,
            "recommended_brand_style": _suggest_brand_style(color_palette),
            "layout_suggestions": _analyze_layout_structure(width, height),
        }

        return json.dumps(analysis_result, indent=2)

    except Exception as e:
        return json.dumps({"error": f"Image analysis failed: {str(e)}"})


@mcp.tool()
async def create_brand_guidelines(
    primary_color: str = "#002F6C",
    secondary_color: str = "#FFFFFF",
    accent_color: str = "#DC3545",
    font_family: str = "Calibri",
    brand_voice: str = "professional",
    confidentiality_level: str = "standard",
) -> str:
    """
    Create custom brand guidelines for presentation generation.

    Args:
        primary_color: Primary brand color (hex)
        secondary_color: Secondary brand color (hex)
        accent_color: Accent color for highlights (hex)
        font_family: Primary font family
        brand_voice: Brand voice (professional, friendly, authoritative, innovative)
        confidentiality_level: Content confidentiality level (public, internal, confidential)

    Returns:
        JSON string with brand guidelines configuration
    """
    try:
        # Convert hex colors to RGB
        def hex_to_rgb(hex_color):
            hex_color = hex_color.lstrip("#")
            return tuple(int(hex_color[i : i + 2], 16) for i in (0, 2, 4))

        brand_config = {
            "brand_guidelines": {
                "colors": {
                    "primary": hex_to_rgb(primary_color),
                    "secondary": hex_to_rgb(secondary_color),
                    "accent": hex_to_rgb(accent_color),
                },
                "typography": {
                    "font_family": font_family,
                    "title_size": 28,
                    "body_size": 14,
                    "caption_size": 10,
                },
                "brand_voice": brand_voice,
                "confidentiality_level": confidentiality_level,
            },
            "voice_guidelines": _get_voice_guidelines(brand_voice),
            "confidentiality_rules": _get_confidentiality_rules(confidentiality_level),
            "created_at": datetime.now().isoformat(),
        }

        return json.dumps(brand_config, indent=2)

    except Exception as e:
        return json.dumps({"error": f"Brand guidelines creation failed: {str(e)}"})


@mcp.tool()
async def extract_pdf_for_presentation(pdf_path: str, focus_areas: str = "") -> str:
    """
    Extract and structure content from a PDF specifically for presentation generation.
    Integrates with the PDF analyzer MCP server functionality.

    Args:
        pdf_path: Path to the PDF file
        focus_areas: Comma-separated focus areas (e.g., "executive summary, recommendations, data")

    Returns:
        JSON string with structured content ready for presentation
    """
    try:
        if not os.path.exists(pdf_path):
            return json.dumps({"error": f"PDF file not found: {pdf_path}"})

        # Extract text using PyMuPDF (same as PDF analyzer)
        doc = fitz.open(pdf_path)
        full_text = ""
        page_texts = []

        for page_num, page in enumerate(doc):
            page_text = page.get_text()
            page_texts.append({"page": page_num + 1, "text": page_text})
            full_text += page_text + "\n"

        doc.close()

        if not full_text.strip():
            return json.dumps({"error": "No text could be extracted from the PDF"})

        # Clean and process text
        cleaned_text = _clean_text_for_presentation(full_text)

        # Extract keywords using TF-IDF (similar to PDF analyzer)
        keywords = _extract_presentation_keywords(cleaned_text)

        # Structure content for presentation
        sections = ContentProcessor._identify_sections(cleaned_text)
        key_points = ContentProcessor._extract_key_points(cleaned_text)

        # Extract charts/data references
        chart_indicators = _identify_chart_content(cleaned_text)

        # Focus on specific areas if specified
        if focus_areas:
            focus_list = [area.strip().lower() for area in focus_areas.split(",")]
            sections = _filter_sections_by_focus(sections, focus_list)
            key_points = _filter_points_by_focus(key_points, focus_list)

        result = {
            "pdf_path": pdf_path,
            "extraction_summary": {
                "total_pages": len(page_texts),
                "text_length": len(cleaned_text),
                "sections_found": len(sections),
                "key_points_found": len(key_points),
            },
            "keywords": keywords,
            "sections": sections[:10],  # Limit to top 10 sections
            "key_points": key_points[:15],  # Limit to top 15 points
            "chart_indicators": chart_indicators,
            "focus_areas": focus_areas.split(",") if focus_areas else [],
            "presentation_ready": True,
            "extracted_at": datetime.now().isoformat(),
        }

        return json.dumps(result, indent=2)

    except Exception as e:
        return json.dumps({"error": f"PDF extraction failed: {str(e)}"})


@mcp.tool()
async def validate_presentation_content(
    content: str, brand_voice: str = "professional"
) -> str:
    """
    Validate presentation content for brand voice compliance, confidentiality, and quality.

    Args:
        content: Text content to validate
        brand_voice: Target brand voice (professional, friendly, authoritative, innovative)

    Returns:
        JSON string with validation results and recommendations
    """
    try:
        # Safety and confidentiality check
        safety_check = ContentSafeguards.check_content_safety(content)

        # Brand voice analysis
        voice_analysis = _analyze_brand_voice_compliance(content, brand_voice)

        # Content quality check
        quality_check = _assess_content_quality(content)

        # Hallucination risk assessment
        hallucination_risk = _assess_hallucination_risk(content)

        validation_result = {
            "content_length": len(content),
            "safety_check": safety_check,
            "brand_voice_analysis": voice_analysis,
            "quality_assessment": quality_check,
            "hallucination_risk": hallucination_risk,
            "overall_score": _calculate_overall_score(
                safety_check, voice_analysis, quality_check, hallucination_risk
            ),
            "recommendations": _generate_content_recommendations(
                safety_check, voice_analysis, quality_check
            ),
            "validated_at": datetime.now().isoformat(),
        }

        return json.dumps(validation_result, indent=2)

    except Exception as e:
        return json.dumps({"error": f"Content validation failed: {str(e)}"})


# Helper Functions


def _generate_slide_structure(
    presentation_type: str,
    title: str,
    content_source: str,
    pdf_content: Dict,
    slide_count: int,
) -> Dict:
    """Generate slide structure based on presentation type"""

    if presentation_type == "proposal":
        return _generate_proposal_structure(
            title, content_source, pdf_content, slide_count
        )
    elif presentation_type == "brand-identity":
        return _generate_brand_identity_structure(title, content_source, slide_count)
    elif presentation_type == "executive-summary":
        return _generate_executive_summary_structure(
            title, content_source, pdf_content, slide_count
        )
    else:
        return _generate_standard_structure(title, content_source, slide_count)


def _generate_proposal_structure(
    title: str, content_source: str, pdf_content: Dict, slide_count: int
) -> Dict:
    """Generate McKinsey-style proposal structure"""

    structure = {
        "subtitle": "Executive Proposal",
        "author": "Strategic Consulting Team",
        "slides": [],
    }

    # Standard proposal flow
    proposal_flow = [
        {"title": "Executive Summary", "type": "summary"},
        {"title": "Problem Statement", "type": "problem"},
        {"title": "Current State Analysis", "type": "analysis"},
        {"title": "Proposed Solution", "type": "solution"},
        {"title": "Implementation Roadmap", "type": "roadmap"},
        {"title": "Expected Benefits", "type": "benefits"},
        {"title": "Risk Assessment", "type": "risks"},
        {"title": "Investment Required", "type": "investment"},
        {"title": "Recommendations", "type": "recommendations"},
        {"title": "Next Steps", "type": "next_steps"},
    ]

    # Use PDF content if available
    if pdf_content and "sections" in pdf_content:
        for i, section in enumerate(pdf_content["sections"][: slide_count - 1]):
            slide = {
                "title": section["title"],
                "bullets": section["content"][:6],  # Limit bullets
                "type": "standard",
            }
            structure["slides"].append(slide)
    else:
        # Generate from content source or use template
        for i, slide_template in enumerate(proposal_flow[: slide_count - 1]):
            bullets = _generate_bullets_for_slide_type(
                slide_template["type"], content_source
            )
            slide = {
                "title": slide_template["title"],
                "bullets": bullets,
                "type": slide_template["type"],
            }
            structure["slides"].append(slide)

    return structure


def _generate_brand_identity_structure(
    title: str, content_source: str, slide_count: int
) -> Dict:
    """Generate brand identity presentation structure"""

    structure = {
        "subtitle": "Brand Identity Guidelines",
        "author": "Brand Strategy Team",
        "slides": [
            {
                "title": "Brand Overview",
                "bullets": [
                    "Mission and vision alignment",
                    "Core brand values and principles",
                    "Brand positioning statement",
                    "Target audience definition",
                ],
            },
            {
                "title": "Visual Identity System",
                "bullets": [
                    "Logo usage guidelines and variations",
                    "Color palette specifications",
                    "Typography hierarchy and fonts",
                    "Imagery style and guidelines",
                ],
            },
            {
                "title": "Brand Voice & Messaging",
                "bullets": [
                    "Tone of voice characteristics",
                    "Key messaging framework",
                    "Communication style guide",
                    "Brand personality traits",
                ],
            },
            {
                "title": "Application Guidelines",
                "bullets": [
                    "Digital asset specifications",
                    "Print material standards",
                    "Presentation templates",
                    "Social media guidelines",
                ],
            },
        ],
    }

    return structure


def _generate_executive_summary_structure(
    title: str, content_source: str, pdf_content: Dict, slide_count: int
) -> Dict:
    """Generate executive summary structure"""

    structure = {
        "subtitle": "Executive Summary",
        "author": "Leadership Team",
        "slides": [],
    }

    if pdf_content and "key_points" in pdf_content:
        # Group key points into slides
        points_per_slide = max(1, len(pdf_content["key_points"]) // (slide_count - 1))

        for i in range(0, len(pdf_content["key_points"]), points_per_slide):
            slide_points = pdf_content["key_points"][i : i + points_per_slide]
            slide = {
                "title": f"Key Insights {i // points_per_slide + 1}",
                "bullets": slide_points,
                "type": "standard",
            }
            structure["slides"].append(slide)
    else:
        # Default executive summary structure
        exec_slides = [
            {
                "title": "Key Findings",
                "bullets": [
                    "Primary insights from analysis",
                    "Critical success factors identified",
                    "Market opportunity assessment",
                    "Competitive landscape overview",
                ],
            },
            {
                "title": "Strategic Recommendations",
                "bullets": [
                    "Immediate action items",
                    "Medium-term strategic initiatives",
                    "Long-term vision alignment",
                    "Resource allocation priorities",
                ],
            },
            {
                "title": "Expected Outcomes",
                "bullets": [
                    "Quantifiable business impact",
                    "Timeline for value realization",
                    "Success metrics and KPIs",
                    "Risk mitigation strategies",
                ],
            },
        ]
        structure["slides"] = exec_slides

    return structure


def _generate_standard_structure(
    title: str, content_source: str, slide_count: int
) -> Dict:
    """Generate standard presentation structure"""

    structure = {"subtitle": "Professional Presentation", "author": "", "slides": []}

    # Create slides based on content source
    if content_source:
        # Simple content splitting
        content_parts = content_source.split("\n\n")
        for i, part in enumerate(content_parts[: slide_count - 1]):
            if part.strip():
                slide = {
                    "title": f"Key Point {i + 1}",
                    "bullets": [part.strip()],
                    "type": "standard",
                }
                structure["slides"].append(slide)

    return structure


def _generate_bullets_for_slide_type(
    slide_type: str, content_source: str = ""
) -> List[str]:
    """Generate appropriate bullet points for different slide types"""

    bullet_templates = {
        "summary": [
            "Executive overview of key findings",
            "Strategic recommendations summary",
            "Expected business impact",
            "Implementation timeline",
        ],
        "problem": [
            "Current challenges identified",
            "Root cause analysis",
            "Impact on business operations",
            "Urgency for resolution",
        ],
        "solution": [
            "Proposed approach overview",
            "Key solution components",
            "Implementation methodology",
            "Success criteria definition",
        ],
        "benefits": [
            "Quantifiable business value",
            "Strategic advantages",
            "Operational improvements",
            "Long-term impact",
        ],
        "risks": [
            "Identified risk factors",
            "Mitigation strategies",
            "Contingency planning",
            "Monitoring approach",
        ],
    }

    return bullet_templates.get(
        slide_type,
        [
            "Key insight or finding",
            "Supporting evidence",
            "Strategic implication",
            "Recommended action",
        ],
    )


def _suggest_brand_style(color_palette: List[Tuple[int, int, int]]) -> str:
    """Suggest brand style based on color analysis"""

    # Analyze dominant colors
    primary_color = color_palette[0] if color_palette else (255, 255, 255)

    # Simple color-based style suggestion
    r, g, b = primary_color

    if r < 50 and g < 50 and b > 100:  # Blue dominant
        return "mckinsey"
    elif r > 200 and g < 100 and b < 100:  # Red dominant
        return "modern"
    else:
        return "corporate"


def _analyze_layout_structure(width: int, height: int) -> Dict:
    """Analyze layout structure from dimensions"""

    aspect_ratio = width / height

    suggestions = {
        "slide_format": "16:9" if 1.7 <= aspect_ratio <= 1.8 else "4:3",
        "layout_type": "wide" if aspect_ratio > 1.6 else "standard",
        "recommended_margins": "0.5 inches" if aspect_ratio > 1.6 else "0.4 inches",
        "text_columns": 2 if width > 1200 else 1,
    }

    return suggestions


def _get_voice_guidelines(brand_voice: str) -> Dict:
    """Get voice guidelines for different brand voices"""

    guidelines = {
        "professional": {
            "tone": "Formal, authoritative, confident",
            "vocabulary": "Technical terms appropriate, avoid jargon",
            "sentence_structure": "Clear, concise, direct",
            "avoid": "Casual language, slang, overly complex sentences",
        },
        "friendly": {
            "tone": "Warm, approachable, conversational",
            "vocabulary": "Accessible language, minimal jargon",
            "sentence_structure": "Varied, engaging, personable",
            "avoid": "Overly formal language, corporate speak",
        },
        "authoritative": {
            "tone": "Expert, commanding, decisive",
            "vocabulary": "Industry-specific terms, precise language",
            "sentence_structure": "Strong, declarative statements",
            "avoid": "Uncertain language, qualifiers, hedging",
        },
        "innovative": {
            "tone": "Forward-thinking, creative, dynamic",
            "vocabulary": "Modern terminology, future-focused",
            "sentence_structure": "Energetic, varied, compelling",
            "avoid": "Traditional corporate language, conservative phrasing",
        },
    }

    return guidelines.get(brand_voice, guidelines["professional"])


def _get_confidentiality_rules(confidentiality_level: str) -> Dict:
    """Get confidentiality rules for different levels"""

    rules = {
        "public": {
            "restrictions": ["No internal company information", "No sensitive data"],
            "allowed": ["General industry information", "Public company information"],
            "watermark": "Public Information",
        },
        "internal": {
            "restrictions": ["No confidential data", "No personal information"],
            "allowed": ["Internal processes", "Company strategy"],
            "watermark": "Internal Use Only",
        },
        "confidential": {
            "restrictions": ["Highly sensitive", "Need-to-know basis"],
            "allowed": ["Strategic information", "Financial data"],
            "watermark": "Confidential",
        },
    }

    return rules.get(confidentiality_level, rules["internal"])


def _clean_text_for_presentation(text: str) -> str:
    """Clean text specifically for presentation use"""

    # Remove excessive whitespace and newlines
    text = re.sub(r"\s+", " ", text)

    # Remove special characters but keep basic punctuation
    text = re.sub(r"[^\w\s.,;:!?()-]", " ", text)

    # Remove URLs and email addresses
    text = re.sub(
        r"http[s]?://(?:[a-zA-Z]|[0-9]|[$-_@.&+]|[!*\\(\\),]|(?:%[0-9a-fA-F][0-9a-fA-F]))+",
        "",
        text,
    )
    text = re.sub(r"\S+@\S+", "", text)

    # Remove page numbers and headers/footers
    text = re.sub(r"\bPage\s+\d+\b", "", text, flags=re.IGNORECASE)
    text = re.sub(r"\b\d+\s+of\s+\d+\b", "", text)

    return text.strip()


def _extract_presentation_keywords(text: str, top_n: int = 15) -> List[str]:
    """Extract keywords specifically for presentations"""

    try:
        # Use TF-IDF for keyword extraction
        sentences = text.split(".")

        if len(sentences) < 2:
            # Fallback to simple word frequency
            words = re.findall(r"\b\w+\b", text.lower())
            word_freq = Counter(words)
            return [word for word, _ in word_freq.most_common(top_n) if len(word) > 3]

        vectorizer = TfidfVectorizer(
            max_features=top_n * 2,
            ngram_range=(1, 2),
            stop_words="english",
            min_df=1,
            max_df=0.95,
        )

        tfidf_matrix = vectorizer.fit_transform(sentences)
        feature_names = vectorizer.get_feature_names_out()

        # Get average TF-IDF scores
        mean_scores = np.mean(tfidf_matrix.toarray(), axis=0)
        top_indices = mean_scores.argsort()[-top_n:][::-1]

        keywords = [feature_names[i] for i in top_indices if mean_scores[i] > 0]
        return keywords

    except Exception as e:
        print(f"Keyword extraction failed: {e}", file=sys.stderr)
        # Simple fallback
        words = re.findall(r"\b\w+\b", text.lower())
        word_freq = Counter(words)
        return [word for word, _ in word_freq.most_common(top_n) if len(word) > 3]


def _identify_chart_content(text: str) -> List[Dict]:
    """Identify potential chart content in text"""

    chart_indicators = []

    # Look for numerical data patterns
    number_patterns = [
        r"\b\d+%\b",  # Percentages
        r"\$\d+[,\d]*\b",  # Currency
        r"\b\d+[,\d]*\s*(?:million|billion|thousand)\b",  # Large numbers
        r"\b\d+\.\d+\b",  # Decimals
    ]

    for pattern in number_patterns:
        matches = re.findall(pattern, text, re.IGNORECASE)
        if matches:
            chart_indicators.append(
                {
                    "type": "numerical_data",
                    "pattern": pattern,
                    "matches": matches[:5],  # Limit matches
                }
            )

    # Look for comparison language
    comparison_words = [
        "compared to",
        "versus",
        "higher than",
        "lower than",
        "increased",
        "decreased",
    ]
    for word in comparison_words:
        if word in text.lower():
            chart_indicators.append(
                {"type": "comparison", "indicator": word, "suitable_for": "bar_chart"}
            )

    return chart_indicators


def _filter_sections_by_focus(
    sections: List[Dict], focus_areas: List[str]
) -> List[Dict]:
    """Filter sections based on focus areas"""

    filtered_sections = []

    for section in sections:
        section_text = f"{section['title']} {' '.join(section['content'])}".lower()

        for focus_area in focus_areas:
            if focus_area in section_text:
                filtered_sections.append(section)
                break

    return filtered_sections or sections[:5]  # Return original if no matches


def _filter_points_by_focus(key_points: List[str], focus_areas: List[str]) -> List[str]:
    """Filter key points based on focus areas"""

    filtered_points = []

    for point in key_points:
        point_lower = point.lower()

        for focus_area in focus_areas:
            if focus_area in point_lower:
                filtered_points.append(point)
                break

    return filtered_points or key_points[:10]  # Return original if no matches


def _analyze_brand_voice_compliance(content: str, brand_voice: str) -> Dict:
    """Analyze content for brand voice compliance"""

    voice_indicators = {
        "professional": ["strategic", "analysis", "recommend", "implement", "optimize"],
        "friendly": ["we", "you", "together", "help", "support"],
        "authoritative": ["must", "will", "proven", "established", "industry-leading"],
        "innovative": [
            "breakthrough",
            "cutting-edge",
            "revolutionary",
            "transform",
            "disrupt",
        ],
    }

    content_lower = content.lower()
    target_indicators = voice_indicators.get(
        brand_voice, voice_indicators["professional"]
    )

    matches = sum(1 for indicator in target_indicators if indicator in content_lower)
    compliance_score = min(100, (matches / len(target_indicators)) * 100)

    return {
        "target_voice": brand_voice,
        "compliance_score": compliance_score,
        "matches_found": matches,
        "total_indicators": len(target_indicators),
        "status": "compliant" if compliance_score >= 60 else "needs_improvement",
    }


def _assess_content_quality(content: str) -> Dict:
    """Assess overall content quality"""

    # Basic quality metrics
    word_count = len(content.split())
    sentence_count = len([s for s in content.split(".") if s.strip()])
    avg_sentence_length = word_count / max(1, sentence_count)

    # Check for common quality issues
    issues = []

    if avg_sentence_length > 30:
        issues.append("Sentences may be too long for presentations")

    if word_count < 50:
        issues.append("Content may be too brief")

    if content.count("!") > word_count * 0.05:
        issues.append("Excessive use of exclamation marks")

    quality_score = max(0, 100 - len(issues) * 20)

    return {
        "word_count": word_count,
        "sentence_count": sentence_count,
        "avg_sentence_length": round(avg_sentence_length, 1),
        "quality_score": quality_score,
        "issues": issues,
        "status": "good"
        if quality_score >= 80
        else "acceptable"
        if quality_score >= 60
        else "needs_improvement",
    }


def _assess_hallucination_risk(content: str) -> Dict:
    """Assess risk of hallucinated or false information"""

    risk_indicators = [
        r"\b(?:definitely|certainly|absolutely|always|never)\b",  # Absolute statements
        r"\b(?:studies show|research indicates|experts say)\b",  # Unsourced claims
        r"\b\d{4}\b",  # Specific years that might be incorrect
        r"\b\d+%\b",  # Specific percentages
    ]

    risk_count = 0
    risk_details = []

    for pattern in risk_indicators:
        matches = re.findall(pattern, content, re.IGNORECASE)
        if matches:
            risk_count += len(matches)
            risk_details.append(
                {"type": pattern, "count": len(matches), "examples": matches[:2]}
            )

    total_words = len(content.split())
    risk_ratio = risk_count / max(1, total_words) * 100

    risk_level = "high" if risk_ratio > 5 else "medium" if risk_ratio > 2 else "low"

    return {
        "risk_level": risk_level,
        "risk_ratio": round(risk_ratio, 2),
        "total_indicators": risk_count,
        "risk_details": risk_details,
        "recommendations": _get_hallucination_recommendations(risk_level),
    }


def _get_hallucination_recommendations(risk_level: str) -> List[str]:
    """Get recommendations based on hallucination risk level"""

    recommendations = {
        "high": [
            "Verify all statistical claims and data points",
            "Add source citations for factual statements",
            "Reduce absolute language (always, never, definitely)",
            "Review specific dates and percentages for accuracy",
        ],
        "medium": [
            "Consider adding qualifying language for uncertain claims",
            "Verify key statistics and dates",
            "Add disclaimers for projections and estimates",
        ],
        "low": [
            "Content appears to have low hallucination risk",
            "Consider adding sources for enhanced credibility",
        ],
    }

    return recommendations.get(risk_level, recommendations["low"])


def _calculate_overall_score(
    safety_check: Dict,
    voice_analysis: Dict,
    quality_check: Dict,
    hallucination_risk: Dict,
) -> Dict:
    """Calculate overall content score"""

    # Weight different factors
    safety_weight = 0.4
    voice_weight = 0.2
    quality_weight = 0.2
    hallucination_weight = 0.2

    safety_score = 100 if safety_check["safe"] else 0
    voice_score = voice_analysis["compliance_score"]
    quality_score = quality_check["quality_score"]

    # Convert hallucination risk to score
    hallucination_scores = {"low": 100, "medium": 70, "high": 30}
    hallucination_score = hallucination_scores.get(hallucination_risk["risk_level"], 50)

    overall_score = (
        safety_score * safety_weight
        + voice_score * voice_weight
        + quality_score * quality_weight
        + hallucination_score * hallucination_weight
    )

    return {
        "overall_score": round(overall_score, 1),
        "grade": "A"
        if overall_score >= 90
        else "B"
        if overall_score >= 80
        else "C"
        if overall_score >= 70
        else "D",
        "component_scores": {
            "safety": safety_score,
            "voice_compliance": voice_score,
            "quality": quality_score,
            "hallucination_risk": hallucination_score,
        },
    }


def _generate_content_recommendations(
    safety_check: Dict, voice_analysis: Dict, quality_check: Dict
) -> List[str]:
    """Generate actionable content recommendations"""

    recommendations = []

    if not safety_check["safe"]:
        recommendations.append(
            "Address confidentiality and safety concerns before using content"
        )

    if voice_analysis["compliance_score"] < 70:
        recommendations.append(
            f"Adjust content to better match {voice_analysis['target_voice']} brand voice"
        )

    if quality_check["quality_score"] < 80:
        recommendations.extend(quality_check["issues"])

    if not recommendations:
        recommendations.append("Content meets quality standards for presentation use")

    return recommendations


if __name__ == "__main__":
    print("PowerPoint Generator MCP Server starting...", file=sys.stderr)
    try:
        mcp.run(transport="stdio")
    except Exception as e:
        print(f"Server error: {e}", file=sys.stderr)
        raise
